#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def add_text_to_slide(slide, text, font_size, top_position, bold=False, color=None):
    """在幻灯片上添加居中文本"""
    left = Inches(0.5)
    top = Inches(top_position)
    width = Inches(9)
    height = Inches(1.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = 'Helvetica Neue'

    if color:
        p.font.color.rgb = color

    return text_box

def create_apple_style_slide(prs, background_color, main_text, main_size,
                             sub_text=None, sub_size=None):
    """创建苹果风格的幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color

    # 主文本
    if sub_text:
        add_text_to_slide(slide, main_text, main_size, 2.5, bold=True,
                         color=RGBColor(255, 255, 255) if background_color == RGBColor(0, 0, 0) else RGBColor(29, 29, 31))
        add_text_to_slide(slide, sub_text, sub_size, 4.2, bold=False,
                         color=RGBColor(200, 200, 200) if background_color == RGBColor(0, 0, 0) else RGBColor(134, 134, 139))
    else:
        add_text_to_slide(slide, main_text, main_size, 3, bold=True,
                         color=RGBColor(255, 255, 255) if background_color == RGBColor(0, 0, 0) else RGBColor(29, 29, 31))

def create_macbook_presentation():
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 定义颜色方案
    black = RGBColor(0, 0, 0)
    white = RGBColor(255, 255, 255)
    light_gray = RGBColor(246, 246, 246)

    # 幻灯片1：封面 - 黑色背景
    create_apple_style_slide(prs, black, "MacBook Pro", 88)

    # 幻灯片2：标语
    create_apple_style_slide(prs, black, "专业人士的\n终极工具", 72)

    # 幻灯片3：芯片介绍
    create_apple_style_slide(prs, white, "M3 Pro\nM3 Max", 80)

    # 幻灯片4：性能数字
    create_apple_style_slide(prs, black, "最多 40 核 GPU", 64,
                            "图形处理性能提升\n最高达 6 倍", 36)

    # 幻灯片5：内存
    create_apple_style_slide(prs, white, "最高 128GB\n统一内存", 68)

    # 幻灯片6：显示屏
    create_apple_style_slide(prs, black, "Liquid Retina XDR", 66,
                            "1600 尼特峰值亮度", 40)

    # 幻灯片7：ProMotion
    create_apple_style_slide(prs, white, "ProMotion\n自适应刷新率", 64,
                            "最高 120Hz", 44)

    # 幻灯片8：续航
    create_apple_style_slide(prs, black, "22 小时", 96,
                            "电池使用时间", 40)

    # 幻灯片9：端口
    create_apple_style_slide(prs, white, "Thunderbolt 4\nHDMI\nSDXC\nMagSafe 3", 52)

    # 幻灯片10：尺寸选择
    create_apple_style_slide(prs, black, "14.2″", 72,
                            "或", 40)

    # 幻灯片11：16英寸
    create_apple_style_slide(prs, black, "16.2″", 72)

    # 幻灯片12：适用场景
    create_apple_style_slide(prs, white, "视频剪辑\n3D 渲染\n软件开发\n音乐制作", 54)

    # 幻灯片13：视频能力
    create_apple_style_slide(prs, black, "8K 视频\n流畅剪辑", 68)

    # 幻灯片14：编译速度
    create_apple_style_slide(prs, white, "编译速度\n提升 2.5 倍", 64)

    # 幻灯片15：设计
    create_apple_style_slide(prs, black, "精致\n强大", 80)

    # 幻灯片16：材质
    create_apple_style_slide(prs, white, "铝金属一体式\n机身", 64)

    # 幻灯片17：颜色
    create_apple_style_slide(prs, light_gray, "深空灰\n银色", 68)

    # 幻灯片18：键盘
    create_apple_style_slide(prs, black, "妙控键盘", 72,
                            "配备触控 ID", 40)

    # 幻灯片19：音频
    create_apple_style_slide(prs, white, "六扬声器\n音响系统", 64,
                            "空间音频支持", 40)

    # 幻灯片20：价格起点
    create_apple_style_slide(prs, black, "¥15,999 起", 68,
                            "14 英寸 MacBook Pro", 40)

    # 幻灯片21：总结
    create_apple_style_slide(prs, white, "性能\n不设限", 84)

    # 幻灯片22：结束语
    create_apple_style_slide(prs, black, "MacBook Pro", 76,
                            "现已发售", 44)

    # 保存演示文稿
    prs.save('MacBook_Pro介绍_苹果风格.pptx')
    print("✓ 苹果风格PowerPoint文件已成功创建：MacBook_Pro介绍_苹果风格.pptx")
    print("✓ 共 22 页，采用极简设计风格")
    print("✓ 黑白配色，大字体，每页一个核心信息")

if __name__ == "__main__":
    create_macbook_presentation()
